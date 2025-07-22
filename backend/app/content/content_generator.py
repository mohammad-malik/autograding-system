import io
import json
import uuid
from typing import Dict, List, Optional, Tuple, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

from ..models import TaskStatus, Quiz, QuestionType
from ..services import LLMClient, StorageClient
from .pdf_processor import PDFProcessor


class ContentGenerator:
    """Content generator for slides and quizzes."""

    @staticmethod
    async def generate_slides(
        prompt: str, textbook_id: str, num_slides: int, user_id: str
    ) -> Tuple[str, str, TaskStatus]:
        """Generate PowerPoint slides."""
        # Search for relevant content
        search_results = await PDFProcessor.search_textbook_content(
            textbook_id, prompt, top_k=20
        )
        
        # Extract text from search results
        context = "\n\n".join(
            [result["metadata"]["text"] for result in search_results]
        )
        
        # Generate slide content using LLM
        slide_data = await LLMClient.generate_slides(prompt, context, num_slides)
        
        # Create PowerPoint presentation
        presentation = Presentation()
        
        # Add title slide
        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = slide_data["slides"][0]["title"]
        if hasattr(title_slide.shapes, "placeholders") and len(title_slide.shapes.placeholders) > 1:
            subtitle = title_slide.shapes.placeholders[1]
            subtitle.text = "Generated with AI"
        
        # Add content slides
        for slide_info in slide_data["slides"][1:]:
            content_slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(content_slide_layout)
            
            # Set title
            slide.shapes.title.text = slide_info["title"]
            
            # Add bullet points
            content_placeholder = slide.shapes.placeholders[1]
            text_frame = content_placeholder.text_frame
            
            for i, bullet_point in enumerate(slide_info["content"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0
            
            # Add speaker notes
            if "notes" in slide_info and slide_info["notes"]:
                slide.notes_slide.notes_text_frame.text = slide_info["notes"]
        
        # Save presentation to memory
        pptx_stream = io.BytesIO()
        presentation.save(pptx_stream)
        pptx_stream.seek(0)
        
        # Upload to storage
        file_path, file_url = await StorageClient.upload_generated_file(
            pptx_stream.getvalue(),
            f"{prompt[:30].replace(' ', '_')}.pptx",
            user_id,
            StorageClient.BUCKET_SLIDES,
        )
        
        return file_path, file_url, TaskStatus.COMPLETED

    @staticmethod
    async def generate_quiz(
        prompt: str,
        textbook_id: str,
        question_types: List[str],
        difficulty: str,
        user_id: str,
    ) -> Tuple[Quiz, str, TaskStatus]:
        """Generate quiz document."""
        # Search for relevant content
        search_results = await PDFProcessor.search_textbook_content(
            textbook_id, prompt, top_k=20
        )
        
        # Extract text from search results
        context = "\n\n".join(
            [result["metadata"]["text"] for result in search_results]
        )
        
        # Generate quiz content using LLM
        quiz_data = await LLMClient.generate_quiz(
            prompt, context, question_types, difficulty
        )
        
        # Create Word document
        document = Document()
        
        # Add title
        title = document.add_heading(quiz_data["title"], level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add description
        if "description" in quiz_data and quiz_data["description"]:
            description = document.add_paragraph(quiz_data["description"])
            description.alignment = WD_ALIGN_PARAGRAPH.CENTER
            document.add_paragraph()  # Add space
        
        # Add questions
        for i, question in enumerate(quiz_data["questions"]):
            # Add question text
            q_text = document.add_paragraph(f"{i+1}. {question['text']}")
            q_text.style = "List Number"
            
            # Add options or space for answer based on question type
            if question["type"] == "multiple_choice" and "options" in question:
                for j, option in enumerate(question["options"]):
                    option_text = document.add_paragraph(f"    {chr(97+j)}) {option}")
                    option_text.style = "List Bullet"
            elif question["type"] == "true_false":
                document.add_paragraph("    True / False")
            else:  # short_answer or essay
                document.add_paragraph("    Answer: ________________________________")
                document.add_paragraph()
        
        # Save document to memory
        docx_stream = io.BytesIO()
        document.save(docx_stream)
        docx_stream.seek(0)
        
        # Upload to storage
        file_path, file_url = await StorageClient.upload_generated_file(
            docx_stream.getvalue(),
            f"{prompt[:30].replace(' ', '_')}_quiz.docx",
            user_id,
            StorageClient.BUCKET_QUIZZES,
        )
        
        # Create quiz record
        quiz_id = str(uuid.uuid4())
        quiz = Quiz(
            id=quiz_id,
            title=quiz_data["title"],
            description=quiz_data.get("description", ""),
            file_path=file_path,
            textbook_id=textbook_id,
            created_by=user_id,
            difficulty=difficulty,
        )
        
        return quiz, file_url, TaskStatus.COMPLETED 